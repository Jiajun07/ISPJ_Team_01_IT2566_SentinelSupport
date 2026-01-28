import os
import zipfile
import json
from io import StringIO
from typing import List, Dict, Optional, Set
from werkzeug.datastructures import FileStorage
from docx import Document
from pptx import Presentation
from pandas import read_excel
from DLPScannerModules.OCRProcessor import performOCRandScan

try:
    import pypdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

class FileProcessor:
    def __init__(self, config_path: Optional[str] = None):
        self.config_path = config_path or os.path.join(os.path.dirname(__file__), "config", "supportedfiles.json")
        self.max_file_size = 10 * 1024 * 1024  # 10 MB
        self.supported_extensions = self.loadSupportedExtensions()
        self.PDF_AVAILABLE = PDF_AVAILABLE
    
    def getFileInfo(self,file):
        if not file or not file.filename:
            return None
        file.seek(0, os.SEEK_END)
        fileSize = file.tell()
        file.seek(0)
        return {
            "filename": file.filename,
            "size": fileSize,
            "extension": os.path.splitext(file.filename)[1].lower()
        }
    
    def loadSupportedExtensions(self) -> Dict[str, Set[str]]:
        if not os.path.exists(self.config_path):
            self._create_default_supported_files_config()
        try:
            with open(self.config_path, 'r', encoding='utf-8') as f:
                config = json.load(f)
            return {
                category: set(extensions)
                for category, extensions in config.items()
            }
        except (FileNotFoundError, json.JSONDecodeError) as e:
            print(f"Error loading supported files config: {e}")
            print("Using default supported file extensions.")
            return self._get_default_supported_extensions()
        
    def _create_default_supported_files_config(self):
        default_extensions = {
            "text_files": ['.txt', '.csv', '.json'],
            "pdf_files": ['.pdf'],
            "archive_files": ['.zip'],
            "office_files": ['.docx', '.xlsx', '.pptx'],
            "image_files": ['.png', '.jpg', '.jpeg', '.bmp', '.gif']
        }
        config_dir = os.path.dirname(self.config_path)
        if not os.path.exists(config_dir):
            os.makedirs(config_dir)
        try:
            with open(self.config_path, 'w', encoding='utf-8') as f:
                json.dump(default_extensions, f, indent=2, ensure_ascii=False)
        except Exception as e:
            raise IOError(f"Failed to create default supported files config: {e}")
    
    def _get_default_supported_extensions(self) -> Dict[str, Set[str]]:
        default_extensions = {
            "text_files": {'.txt', '.csv', '.json'},
            "pdf_files": {'.pdf'},
            "archive_files": {'.zip'},
            "office_files": {'.docx', '.xlsx', '.pptx'},
            "image_files": {'.png', '.jpg', '.jpeg', '.bmp', '.gif'}
        }
        return default_extensions
    
    def reloadConfig(self):
        self.supported_extensions = self.loadSupportedExtensions()
    
    def getAllSupportedExtensions(self) -> Set[str]:
        all_extensions = set()
        for extensions in self.supported_extensions.values():
            all_extensions.update(extensions)
        return all_extensions
    
    def getExtensionsByCategory(self, category: str) -> Set[str]:
        return self.supported_extensions.get(category, set())
    
    def getFileCategory(self, filename: str) -> Optional[str]:
        ext = os.path.splitext(filename)[1].lower()
        for category, extensions in self.supported_extensions.items():
            if ext in extensions:
                return category
        return None

    def getFileExtension(self, filename: str) -> str:
        return os.path.splitext(filename)[1].lower()
    
    def passedProcessing(self, file: FileStorage) -> bool:
        if not file or not file.filename:
            return False
        ext = os.path.splitext(file.filename)[1].lower()
        if ext in self.supported_extensions.get("text_files", set()):
            return True
        if ext in self.supported_extensions.get("pdf_files", set()) and PDF_AVAILABLE:
            return True
        if ext in self.supported_extensions.get("archive_files", set()):
            return True
        if ext in self.supported_extensions.get("office_files", set()):
            return True
        if ext in self.supported_extensions.get("image_files", set()):
            return True
        if ext in self.supported_extensions.get("code_files", set()):
            return True
        return False
    
    def readTextFromFile(self, file: FileStorage) -> str:
        if not file or not file.filename:
            raise ValueError(f"File type not supported: {file.filename}")
        ext = self.getFileExtension(file.filename)
        try:
            if ext in self.supported_extensions.get("text_files", set()):
                return self.readTextFromTextFile(file, ext)
            elif ext in self.supported_extensions.get("pdf_files", set()) and PDF_AVAILABLE:
                return self.readTextFromPDFFile(file)
            elif ext in self.supported_extensions.get("archive_files", set()):
                return self.readTextFromZipFile(file)
            elif ext in self.supported_extensions.get("office_files", set()):
                if ext == '.docx':
                    return self.readTextFromDocxFile(file)
                elif ext == '.xlsx':
                    return self.readTextFromExcelFile(file)
                elif ext == '.pptx':
                    return self.readTextFromPPTXFile(file)
            elif ext in self.supported_extensions.get("image_files", set()):
                return self.readTextFromOCRScanner(file)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            raise ValueError(f"Error reading file: {e}") from e
    
    def readTextFromTextFile(self, file: FileStorage, ext: str) -> str:
        try:
            content = file.read()
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                return content.decode('latin-1', errors='ignore')
        except Exception as e:
            raise ValueError(f"Error reading text file: {e}") from e
    
    def readTextFromDocxFile(self, file: FileStorage) -> str:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx/Microsoft Word Module is not installed.")
        try:
            file.seek(0)
            document = Document(file)
            texts = [para.text for para in document.paragraphs]
            return "\n".join(texts)
        except Exception as e:
            raise ValueError(f"Error reading DOCX file '{file.filename}': {str(e)}") from e
    
    def readTextFromExcelFile(self, file: FileStorage) -> str:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pandas library/Excel Module not installed.")
        try:
            file.seek(0)
            xls = pd.ExcelFile(file)
            texts = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                texts.append(f"--- Sheet: {sheet_name} ---\n")
                texts.append(df.to_string(index=False))
                texts.append("\n")
            return "\n".join(texts)
        except Exception as e:
            raise ValueError(f"Error reading Excel file '{file.filename}': {str(e)}") from e
    
    def readTextFromPPTXFile(self, file: FileStorage) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx library/PowerPoint Module is not installed.")
        
        try:
            file.seek(0)
            presentation = Presentation(file)
            texts = []
            for slide_num, slide in enumerate(presentation.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_texts.append(shape.text)
                texts.append(f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_texts))
            return "\n\n".join(texts)
        except Exception as e:
            raise ValueError(f"Error reading PPTX file '{file.filename}': {str(e)}") from e
    
    def readTextFromPDFFile(self, file: FileStorage) -> str:
        if not PDF_AVAILABLE:
            raise ImportError("pypdf library is not installed.")
        try:
            file.seek(0)
            reader = pypdf.PdfReader(file)
            text = []
            
            for page_num, page in enumerate(reader.pages):
                try:
                    extracted_text = page.extract_text()
                    if extracted_text:
                        text.append(f"--- Page {page_num + 1} ---\n{extracted_text}")
                except Exception as e:
                    text.append(f"--- Page {page_num + 1} (Error reading) ---\nError: {str(e)}")
            if not text:
                return "No text could be extracted from this PDF."
            return "\n\n".join(text)
        except Exception as e:
            raise ValueError(f"Error reading PDF file '{file.filename}': {str(e)}") from e
    
    def readTextFromZipFile(self, file: FileStorage) -> str:
        try:
            texts = []
            with zipfile.ZipFile(file, "r") as zip_ref:
                for file_info in zip_ref.filelist:
                    if file_info.is_dir():
                        continue
                    ext = os.path.splitext(file_info.filename)[1].lower()
                    if ext in self.supported_extensions.get("text_files", set()):
                        try:
                            with zip_ref.open(file_info) as extracted_file:
                                content = extracted_file.read()
                                try:
                                    texts.append(content.decode('utf-8'))
                                except UnicodeDecodeError:
                                    texts.append(content.decode('latin-1', errors='ignore'))

                                texts.append(f"{file_info.filename} (from zip):\n")
                        except Exception as e:
                            texts.append(f"Error reading text file in zip archive: {e}\n")
            return "\n".join(texts)
        except Exception as e:
            raise ValueError(f"Error reading text file in zip archive: {e}") from e
    
    def readTextFromOCRScanner(self, image_path: FileStorage) -> str:
        try:
            ocr_text = performOCRandScan(image_path)
            return ocr_text
        except Exception as e:
            raise ValueError(f"Error performing OCR on image file '{image_path.filename}': {str(e)}") from e